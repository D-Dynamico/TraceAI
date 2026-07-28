"""Multi-format text extraction.

Detects a file's type from its extension and extracts plain text using the
appropriate parser:
  - PDF   -> PyMuPDF (fitz); OCR fallback for scanned/image-only PDFs
  - DOCX  -> python-docx
  - PPTX  -> python-pptx
  - TXT/MD-> raw read
  - Images-> OCR (local Tesseract, then Gemini Vision — see ocr_handler)

Returns an ExtractionResult with the text and metadata about how it was
obtained, so downstream modules (and the UI) can flag low-confidence extractions.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

from ai import degradation
from config import settings
from ingestion import ocr_handler

logger = logging.getLogger(__name__)

# Extension -> logical file type
PDF_EXTS = {".pdf"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}
TEXT_EXTS = {".txt", ".md", ".markdown"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".tif", ".webp"}

SUPPORTED_EXTS = PDF_EXTS | DOCX_EXTS | PPTX_EXTS | TEXT_EXTS | IMAGE_EXTS


@dataclass
class ExtractionResult:
    text: str
    file_type: str          # "pdf" | "docx" | "pptx" | "text" | "image"
    # "native" | "ocr" | "vision" | "native+ocr" | "native+vision".
    # "vision" means Gemini read the pixels because local OCR yielded nothing;
    # it is surfaced on the API so a reviewer can see *which* rung produced the
    # text rather than being told only that some OCR ran.
    method: str
    char_count: int = 0
    # True whenever text came from reading pixels rather than a text layer —
    # local OCR and Vision alike, since both are "this was not machine text".
    used_ocr: bool = False
    warnings: list[str] = field(default_factory=list)
    # Why the text is missing, as a reason code plus `retryable` — the same
    # contract the four Gemini callers degrade through (`ai/degradation.py`).
    # Until now this cause existed only as prose inside `warnings`, so a client
    # could not tell a quota wall (wait, then retry) from a missing Tesseract
    # binary (retrying forever will not help) without parsing a sentence.
    # Categorization has had the structured form since deferred item B;
    # extraction is upstream of it and had nothing.
    #
    # **Invariant: set exactly when the no-text warning is emitted**, so the
    # prose and the code can never disagree about whether extraction failed.
    degraded: degradation.Degradation | None = None

    def __post_init__(self) -> None:
        self.char_count = len(self.text)


class UnsupportedFileError(ValueError):
    """Raised when a file extension has no registered parser."""


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTS


def detect_file_type(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in PDF_EXTS:
        return "pdf"
    if ext in DOCX_EXTS:
        return "docx"
    if ext in PPTX_EXTS:
        return "pptx"
    if ext in TEXT_EXTS:
        return "text"
    if ext in IMAGE_EXTS:
        return "image"
    raise UnsupportedFileError(f"Unsupported file extension: {ext or '(none)'}")


def _extract_pdf(path: Path) -> ExtractionResult:
    import fitz  # PyMuPDF

    warnings: list[str] = []
    native_text = ""
    try:
        with fitz.open(path) as doc:
            native_text = "\n\n".join(page.get_text().strip() for page in doc).strip()
    except Exception as exc:
        warnings.append(f"PyMuPDF failed: {exc}")
        logger.warning("PyMuPDF extraction failed for %s: %s", path, exc)

    # If the PDF has little/no extractable text, it's likely scanned -> OCR.
    if len(native_text) < settings.ocr_char_threshold:
        ocr = ocr_handler.ocr_pdf(path)
        if ocr.text:
            method = f"native+{ocr.method}" if native_text else ocr.method
            combined = (native_text + "\n\n" + ocr.text).strip() if native_text else ocr.text
            return ExtractionResult(combined, "pdf", method, used_ocr=True, warnings=warnings)
        if not native_text:
            warnings.append(_no_text_warning(ocr))
            return ExtractionResult(
                native_text, "pdf", "native", used_ocr=False,
                warnings=warnings, degraded=ocr.degraded,
            )
        # Some native text survived, so extraction did not fail — a short PDF is
        # allowed to be short. No warning, and so (by the invariant) no reason.
        return ExtractionResult(native_text, "pdf", "native", used_ocr=False, warnings=warnings)

    return ExtractionResult(native_text, "pdf", "native", used_ocr=False, warnings=warnings)


def _extract_docx(path: Path) -> ExtractionResult:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    # Include table cell text as well.
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text.strip())
    return ExtractionResult("\n".join(parts).strip(), "docx", "native")


def _extract_pptx(path: Path) -> ExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
    return ExtractionResult("\n\n".join(parts).strip(), "pptx", "native")


def _extract_text(path: Path) -> ExtractionResult:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return ExtractionResult(text, "text", "native")


def _no_text_warning(result: ocr_handler.OcrResult) -> str:
    """Name which rung failed and why, not just that the text is missing.

    The failure this replaces read "OCR produced no text (Tesseract unavailable
    or blank image)" — one sentence covering two unrelated causes with opposite
    fixes. A missing binary is the operator's problem; a quota wall clears on its
    own; a blank scan is nobody's. Each now says so.
    """
    parts: list[str] = ["No text could be extracted."]
    parts.append(
        "Local OCR (Tesseract) is not installed."
        if not result.local_available
        else "Local OCR found no text."
    )
    if result.degraded is not None:
        parts.append(f"Gemini Vision: {result.degraded.message}.")
    return " ".join(parts)


def _extract_image(path: Path) -> ExtractionResult:
    warnings: list[str] = []
    ocr = ocr_handler.ocr_image(path)
    if not ocr.text:
        warnings.append(_no_text_warning(ocr))
    # `method` reports the rung that won; on total failure keep "ocr" so the
    # field still says how the file was *approached*.
    return ExtractionResult(
        ocr.text, "image", ocr.method or "ocr", used_ocr=True,
        warnings=warnings, degraded=ocr.degraded if not ocr.text else None,
    )


_EXTRACTORS = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "pptx": _extract_pptx,
    "text": _extract_text,
    "image": _extract_image,
}


def extract_text(path: Path) -> ExtractionResult:
    """Extract text from a file, dispatching on detected type.

    Raises UnsupportedFileError for unknown extensions.
    """
    file_type = detect_file_type(path)
    extractor = _EXTRACTORS[file_type]
    result = extractor(path)
    logger.info(
        "Extracted %d chars from %s (type=%s, method=%s, ocr=%s)",
        result.char_count, path.name, result.file_type, result.method, result.used_ocr,
    )
    return result
