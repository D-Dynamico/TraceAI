"""Shared fixtures.

Two things every test gets for free:

  - **Isolated storage.** `settings` is a module-level singleton that every
    module imports by reference, so repointing its paths at a tmp directory
    redirects the whole app. Without this, running the suite would write real
    files into `uploads/` and rows into `data/traceai.db`.

  - **A stubbed categorizer.** Tests must not spend Gemini quota or fail when
    the network is down. `categorize()` is replaced with a deterministic stub;
    tests that want the real API opt in via the `live` marker.
"""

from __future__ import annotations

import hashlib
import io
import math

import pytest
from fastapi.testclient import TestClient

import storage
from ai import categorizer, embeddings, vision
from config import settings
from db import database
from models.document import Categorization

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture(autouse=True)
def isolated_storage(tmp_path, monkeypatch):
    """Point uploads, the database, and the vector store at a per-test tmp dir.

    `chroma_dir` matters as much as the others: without it every test that
    ingests a document would write a real Chroma store into the repo's `data/`.
    The embeddings module caches its client/collection in module globals, so
    those are reset per test — otherwise a collection opened against one test's
    tmp dir would leak into the next.
    """
    uploads = tmp_path / "uploads"
    data = tmp_path / "data"
    chroma = tmp_path / "chroma"
    uploads.mkdir()
    data.mkdir()
    chroma.mkdir()

    monkeypatch.setattr(settings, "upload_dir", uploads)
    monkeypatch.setattr(settings, "data_dir", data)
    monkeypatch.setattr(settings, "db_path", data / "traceai.db")
    monkeypatch.setattr(settings, "chroma_dir", chroma)

    # Drop any store cached against a previous test's directory.
    monkeypatch.setattr(embeddings, "_client", None)
    monkeypatch.setattr(embeddings, "_collection", None)

    database.init_db()
    yield tmp_path


@pytest.fixture
def stub_result() -> Categorization:
    """The categorization every stubbed upload returns."""
    return Categorization(
        document_type="certificate",
        category="Certifications",
        title="Test Certificate",
        date="2024-03",
        summary="A stubbed summary.",
        skills=["Python", "SQL"],
        organizations=["Coursera"],
        people=["Ada Lovelace"],
        tags=["testing"],
        confidence=0.9,
    )


@pytest.fixture(autouse=True)
def stub_categorizer(request, monkeypatch, stub_result):
    """Replace the Gemini call with a deterministic stub.

    Opted out of by two markers:
      - `live`   — exercises the real API on purpose.
      - `nostub` — exercises the real `categorize()` logic offline (its
        fallback and error handling). Without this escape hatch such tests
        would silently assert against the stub and pass for the wrong reason.
    """
    if request.node.get_closest_marker("live") or request.node.get_closest_marker("nostub"):
        return

    def _fake(text: str, filename: str = "") -> Categorization:
        return stub_result.model_copy(deep=True)

    # Patched where it is looked up, not just where it is defined.
    monkeypatch.setattr(categorizer, "categorize", _fake)
    import routes.upload as upload_route

    monkeypatch.setattr(upload_route.categorizer, "categorize", _fake)


@pytest.fixture(autouse=True)
def stub_vision(request, monkeypatch):
    """Keep the Gemini Vision rung off the network.

    Extraction gained a Gemini call (`ai/vision.py`, reached from
    `ocr_handler` when local OCR yields nothing), so **any** test that uploads an
    image or a scanned PDF would otherwise spend real free-tier quota. No test
    did when the rung was added — which is exactly why this fixture has to exist
    before one does.

    Unlike `stub_categorizer`, this replaces only `_generate`, the module's sole
    network touch. Everything worth testing in `extract_text` — the config gate,
    the key check, the inline size cap, the sentinel, the exception→reason
    mapping — is code *around* that call and stays live. Stubbing
    `extract_text` wholesale would skip all of it.

    The default is an empty transcript, i.e. "the model saw no text": it leaves
    every pre-existing expectation about images intact, and it fails *closed* —
    a test that needs Vision to succeed has to say so, so no test can pass on
    text that a stub invented. Those tests patch `vision._generate` themselves.
    """
    if request.node.get_closest_marker("live"):
        return

    # Drop any client cached against a previous test's key/config.
    monkeypatch.setattr(vision, "_model", None)
    monkeypatch.setattr(vision, "_generate", lambda data, mime_type: "")


@pytest.fixture(autouse=True)
def stub_embeddings(request, monkeypatch):
    """Replace the sentence-transformer with a deterministic fake.

    The real model is an ~80MB download and slow to load, neither of which the
    offline suite should pay. `embed_texts` is the single choke point every
    embedding flows through, so stubbing it alone covers add, query, and
    reindex. Vectors are a normalized hash of the text: consistent in dimension
    and per-text, enough to exercise the store's mechanics (filtering, dedup,
    hydration). Ranking *quality* needs the real model — those tests opt in with
    the `model` marker, which skips this stub.
    """
    if request.node.get_closest_marker("model"):
        return

    def _fake(texts: list[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        for text in texts:
            digest = hashlib.sha256(text.encode("utf-8")).digest()
            raw = [b / 255.0 for b in digest[:16]]
            norm = math.sqrt(sum(x * x for x in raw)) or 1.0
            vectors.append([x / norm for x in raw])
        return vectors

    monkeypatch.setattr(embeddings, "embed_texts", _fake)
    # The ingest routes now kick off a background model load before embedding.
    # `embed_texts` above never reaches `_get_model`, but prewarm calls it
    # directly — so without this the offline suite would pay for the download
    # this fixture exists to avoid, on a thread where the cost is invisible.
    monkeypatch.setattr(embeddings, "prewarm", lambda: None)


@pytest.fixture
def client() -> TestClient:
    from main import app

    return TestClient(app)


# --- Sample documents ------------------------------------------------------


def make_docx(
    heading: str = "Internship Completion Certificate",
    body: str = "This certifies completion of a Python internship at Acme AI.",
) -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading(heading, level=1)
    doc.add_paragraph(body)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_pptx(title: str = "ML Pipeline Project") -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Built an ML pipeline using scikit-learn."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_image(fmt: str = "PNG", size: tuple[int, int] = (48, 24)) -> bytes:
    """A tiny blank image. Stands in for a scan: what matters to the OCR ladder
    is that no text layer exists, not what the pixels show."""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", size, "white").save(buf, format=fmt)
    return buf.getvalue()


def make_textless_pdf() -> bytes:
    """A valid PDF with a page but no text layer — the scanned-PDF shape.

    PyMuPDF is already a dependency, and a real scan is the same thing to the
    parser: `page.get_text()` returns nothing, so extraction drops to OCR.
    """
    import fitz

    doc = fitz.open()
    doc.new_page()
    data = doc.tobytes()
    doc.close()
    return data


def upload(client: TestClient, name: str, data: bytes, mime: str = "text/plain"):
    return client.post("/api/upload", files={"file": (name, data, mime)})


@pytest.fixture
def stored_doc(client):
    """Upload a docx and return (doc_id, original_bytes, response_json)."""
    original = make_docx()
    resp = upload(client, "cert.docx", original, DOCX_MIME)
    assert resp.status_code == 200, resp.text
    return resp.json()["id"], original, resp.json()
